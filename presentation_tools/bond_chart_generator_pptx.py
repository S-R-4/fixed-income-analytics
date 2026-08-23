from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import numpy as np

# ───────────────────────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = PROJECT_ROOT / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

prs = Presentation() #create presentation
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ───────────────────────────────────────────────────────────────
# 1) Build your data
yields = np.linspace(0, .10, 101)        # 0% → 10% in 0.1% steps
face   = 100
# helper to price a coupon bond
def bond_price(face, coupon_rate, yrs, ytm):
    c = face * coupon_rate
    pv_coupons = sum(c / (1+ytm)**t for t in range(1, yrs+1))
    pv_par     = face/(1+ytm)**yrs
    return pv_coupons + pv_par

prices_30 = [bond_price(face, 0.04625, 30, y) for y in yields]
prices_5  = [bond_price(face, 0.03200,  5, y) for y in yields]
prices_1  = [bond_price(face, 0.00000,  1, y) for y in yields]

# ───────────────────────────────────────────────────────────────
# 2) Create a new blank deck + slide
prs   = Presentation()
blank = prs.slide_layouts[6]           # completely blank
slide = prs.slides.add_slide(blank)

# 3) Add a black header bar & white title
hdr = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(0), Inches(0), 
    prs.slide_width, Inches(0.5)
)
hdr.fill.solid()
hdr.fill.fore_color.rgb = RGBColor(0,0,0)
hdr.line.fill.background()

tf = hdr.text_frame
p  = tf.paragraphs[0]
p.text = "Exhibit 1: Three Bond Prices at Varying Yields-to-Maturity"
p.font.color.rgb = RGBColor(255,255,255)
p.font.bold      = True
p.font.size      = Pt(14)

# ───────────────────────────────────────────────────────────────
# 4) Build the ChartData
chart_data = ChartData()
chart_data.categories = [f"{y*100:.1f}%" for y in yields]  # tick labels
chart_data.add_series("4.625% Coupon, 30-yr Romanian bond", prices_30)
chart_data.add_series("3.2% Coupon, 5-yr BRWA bond",         prices_5)
chart_data.add_series("0% Coupon, 1-yr Australian bond",      prices_1)

# 5) Insert the line chart
x, y, cx, cy = Inches(0.5), Inches(0.7), Inches(9), Inches(4.2)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# 6) Format axes
cat_ax = chart.category_axis
val_ax = chart.value_axis

# percent formatting
for ax in (cat_ax, val_ax):
    ax.has_major_gridlines = False
    ax.tick_labels.number_format = '0%'
    ax.tick_labels.font.size = Pt( 9 )  # tweak as you like
# lighten the gridlines
val_ax.major_gridlines.format.line.width = Pt(0.01)

# 7) Remove the legend (we’ll use inline labels)
chart.has_legend = False

# 8) Style each series
colors = [
    RGBColor(  0,   0, 255),   # blue
    RGBColor(255,   0,   0),   # red
    RGBColor(  0, 150,   0),   # dark green
]
width  = Pt(2)
for ser, col in zip(chart.series, colors):
    ser.format.line.width = width
    ser.format.line.color.rgb = col

# 9) Add inline text‐labels next to the right edge
labels = [
    ("4.625% Coupon, 30-yr\nRomanian bond", Inches(9.1), Inches(1.2)),
    ("3.2% Coupon, 5-yr MAKE bond",         Inches(9.1), Inches(2.0)),
    ("0% Coupon, 1-yr Australian bond",    Inches(9.1), Inches(3.0)),
]
for text, px, py in labels:
    tb = slide.shapes.add_textbox(px, py, Inches(2), Inches(0.6))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(80,80,80)  # a grey tone

# 10) Save
prs.save(OUTPUT_DIR / "Exhibit1_bond_chart.pptx")