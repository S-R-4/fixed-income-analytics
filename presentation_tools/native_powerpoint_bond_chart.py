#Boilerplate: Create (or open a deck and add a blank slide)
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = PROJECT_ROOT / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

prs = Presentation()  # start a new deck
slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout 5 is blank

#3 Compute your Data
import numpy as np

yields = np.arange(0, 10.1, 2)       # 0%, 2%, 4%, … 10%
coupon = 0.04625
maturity = 30

# price = PV of coupons + PV of par
prices = [
    sum([100 * coupon / (1 + y/100)**t for t in range(1, maturity+1)])
    + 100  / (1 + y/100)**maturity
    for y in yields
]

#4 Build a Real PowerPoint Chart
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE

# 4.1 Prepare chart data
chart_data = ChartData()
chart_data.categories = list(yields)           # x‑axis labels
chart_data.add_series('Romanian Bond', prices)

# 4.2 Add the chart to your slide
x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# 4.3 Tweak title, legend, axis titles
chart.chart_title.text_frame.text = "Bond Prices vs. Yield‑to‑Maturity"
chart.has_legend = False

# set axis titles
cat_ax = chart.category_axis
val_ax = chart.value_axis
cat_ax.has_title = True
cat_ax.axis_title.text_frame.text = "Yield-to-Maturity (%)"
val_ax.has_title = True
val_ax.axis_title.text_frame.text = "Price"

# 4.4 Series default formatting (line style, color, etc)
series = chart.series[0]
series.format.line.width = Inches(0.03)
series.marker.symbol = XL_MARKER_STYLE.NONE       # start with no markers

#5 Put a marker at just one point
# choose which index you want marked, e.g. at yields == 4% → index 2
marker_idx = list(yields).index(4)

for i, point in enumerate(series.points):
    if i == marker_idx:
        point.marker.size = 8                    # show a circle
        point.marker.symbol = XL_MARKER_STYLE.CIRCLE
    else:
        point.marker.symbol = XL_MARKER_STYLE.NONE                    # hide all the rest

#6 save your deck
prs.save(OUTPUT_DIR / "bond_yield_chart.pptx")