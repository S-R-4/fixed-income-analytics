from pathlib import Path
from io import BytesIO

#plot matplotlib figure in powerpoint
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = PROJECT_ROOT / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

yields = np.arange(0, 10.1, 2)
coupon = 0.04625
maturity = 30

prices = [
    sum(100 * coupon / (1 + y/100)**t for t in range(1, maturity + 1))
    + 100 / (1 + y/100)**maturity
    for y in yields
]

# … your code to build fig & ax …
fig, ax = plt.subplots(figsize=(8,5))
ax.plot(yields, prices)
ax.set( xlabel="Yield-to-Maturity (%)",
        ylabel="Price",
        title="Bond Prices vs Yield-to-Maturity" )
buf = BytesIO()
fig.savefig(buf, format='png', dpi=150)
buf.seek(0)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))

prs.save(OUTPUT_DIR / "matplotlib_bond_chart.pptx")