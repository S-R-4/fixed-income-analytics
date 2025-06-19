import numpy as np
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# --- Define bond metadata ---
bonds = {
    "Bulgarian Bond": {"coupon": 0.04625, "maturity": 30},
    "Loews Corp Bond": {"coupon": 0.03200, "maturity": 5},
    "Japan Bond": {"coupon": 0.0, "maturity": 1},
}

# --- Bond pricing function ---
def bond_price(face, coupon_rate, maturity, ytm):
    coupon = face * coupon_rate
    price = sum(coupon / (1 + ytm)**t for t in range(1, maturity + 1))
    price += face / (1 + ytm)**maturity
    return price

# --- Generate yields ---
yields = np.linspace(0, 10, 101)  # 0% to 10% in 0.1 steps

# --- Store prices dynamically ---
bond_prices = {}

for name, props in bonds.items():
    coupon = props["coupon"]
    maturity = props["maturity"]

    if coupon == 0.0:
        bond_prices[name] = [100 / (1 + y/100) for y in yields]
    else:
        bond_prices[name] = [bond_price(100, coupon, maturity, y/100) for y in yields]

# --- Plot the chart ---
fig, ax = plt.subplots(figsize=(8, 5))
for name, props in bonds.items():
    coupon = props["coupon"]
    maturity = props["maturity"]

    # Generate formatted label
    if coupon == 0:
        label = f"{name} ({maturity}yr, 0%)"
    else:
        label = f"{name} ({maturity}yr, {coupon * 100:.3f}%)"

    ax.plot(yields, bond_prices[name], label=label)

ax.set(
    xlabel="Yield-to-Maturity (%)",
    ylabel="Bond Price (% of Par)",
    title="Bond Prices vs Yield-to-Maturity"
)
ax.legend()
ax.grid(True)
plt.show()

# --- Save chart to PowerPoint ---
buf = BytesIO()
fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
buf.seek(0)

prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# Remove default placeholders (optional cleanup)
for shp in list(slide.shapes):
    if shp.is_placeholder:
        slide.shapes._spTree.remove(shp._element)

slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(8))

# Save the PPTX
pptx_path = "bond_chart_from_jupyter.pptx"
prs.save(pptx_path)
print(f"Saved to {pptx_path}")
print("Saving to directory:", os.getcwd())