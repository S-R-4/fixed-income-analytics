import numpy as np
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Bond pricing function
def bond_price(face, coupon_rate, years, ytm):
    coupon = face * coupon_rate
    # Present value of coupons + face
    pv_coupons = sum(coupon / (1 + ytm) ** t for t in range(1, years + 1))
    pv_face    = face / (1 + ytm) ** years
    return pv_coupons + pv_face

# identify bonds
bond_1 = "Bulgarian Bond (30yr, 4.625%)"
bond_2 = "Loews Corp Bond (5yr, 3.2%)"
bond_3 = "Japan Bond (1yr, 0%)" #zero-coupon

# generate yields
yields    = np.linspace(0, 10, 101)

# hardcode real data:
prices_b1 = [bond_price(100, 0.04625, 30, y/100) for y in yields]  # 1)
prices_b2  = [bond_price(100, 0.03200,  5, y/100) for y in yields]  # 2)
prices_b3  = [100/(1 + y/100)           for y in yields]            # 3)

# 1) draw the chart
fig, ax = plt.subplots(figsize=(8,5))
ax.plot(yields, prices_b1, label=bond_1)
ax.plot(yields, prices_b2,  label=bond_2)
ax.plot(yields, prices_b3,  label=bond_3)
ax.set(
    xlabel="Yield-to-Maturity (%)",
    ylabel="Bond Price (% of Par)",
    title="Bond Prices vs Yield-to-Maturity"
)
ax.legend()
ax.grid(True)

# 2) capture it
buf = BytesIO()
fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
buf.seek(0)

# 3) paste into PPT
prs   = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
for shp in list(slide.shapes):
    if shp.is_placeholder:
        slide.shapes._spTree.remove(shp._element)
slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(8))

# 4) save
prs.save("bond_chart_from_jupyter.pptx")
print("Saved to bond_chart_from_jupyter.pptx")
print("Saving to:", os.getcwd())

# 5) save to specific folder
#prs.save("reports/bond_chart_from_jupyter.pptx") # change path
#prs.save("C:/Users/.../bond_chart_from_jupyter.pptx") # or fully explicit
